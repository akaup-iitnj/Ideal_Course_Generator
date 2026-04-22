"""
Stage 1: single-file course video pipeline (all steps in this file).

Usage:
  python make_video.py "Lesson title"     # full run -> output/output.mp4
  python make_video.py --test-slides      # slides only (sample script)
  python make_video.py --test-audio       # TTS only (sample script; needs .env)
  python make_video.py --test-compose     # FFmpeg only (needs matching slide_*.png/mp3)
Requires: OPENAI_API_KEY in .env, LibreOffice, ffmpeg on PATH.
Optional: PPTX_TEMPLATE path (see generate_slides) to load a .pptx that defines slide masters; otherwise
a built-in 16:9 layout is used. For a polished "template" look, many teams use Canva, SlidesCarnival, or
Microsoft/Office design ideas, then export; fully automated "premium" decks usually mean either a
custom .potx (opened here as PPTX) or a design service API (Gamma, Beautiful.ai) outside this repo.

Per-slide "image" is a path or URL, or a filename under source_images_dir (set by stage4 for PDF extract).
"""

from __future__ import annotations

import json
import os
import shutil
import subprocess
import sys
import urllib.error
import urllib.request
from pathlib import Path
from typing import Any

# -----------------------------------------------------------------------------


def _openai_client():
    """Shared client; loads .env from this folder. Same key as script + TTS need OpenAI access."""
    from dotenv import load_dotenv
    from openai import OpenAI

    load_dotenv(Path(__file__).resolve().parent / ".env")
    api_key = os.environ.get("OPENAI_API_KEY", "").strip()
    if not api_key:
        raise SystemExit(
            "Missing OPENAI_API_KEY. Set it in stage1/.env (OPENAI_API_KEY=sk-...)."
        )
    return OpenAI(api_key=api_key)


def generate_script(topic: str) -> dict[str, Any]:
    """
    Call GPT-4o to return a script dict: topic + slides (title, bullets, narration per slide).
    Loads OPENAI_API_KEY from a .env file in this folder (or from the environment) via python-dotenv;
    that avoids hardcoding secrets in the script.
    """
    client = _openai_client()

    system = """You write short educational video scripts. Output a single JSON object with this exact structure (no markdown, no extra keys):
{
  "topic": "<same as user topic string>",
  "slides": [
    {
      "slide_num": 1,
      "title": "<short slide title>",
      "bullets": ["<bullet 1>", "<bullet 2>"],
      "narration": "<2-4 sentences the presenter would say for this slide; conversational>"
    }
  ]
}
Rules:
- Produce 6 to 8 slides. Together they should support roughly 2–3 minutes of spoken narration.
- Each slide: exactly 1 title string, 2 to 4 bullets, each bullet one short line.
- narration is only spoken; do not repeat the title verbatim in every narration.
- slide_num must be 1, 2, 3, ... with no gaps.
- Return only valid JSON."""

    response = client.chat.completions.create(
        model="gpt-4o",
        response_format={"type": "json_object"},
        messages=[
            {"role": "system", "content": system},
            {
                "role": "user",
                "content": f"Topic (course lesson title): {topic}",
            },
        ],
    )

    raw = response.choices[0].message.content
    if not raw:
        raise RuntimeError("Empty response from model")

    data: dict[str, Any] = json.loads(raw)
    if "topic" not in data or "slides" not in data:
        raise ValueError("Response must include 'topic' and 'slides'")

    for i, slide in enumerate(data["slides"], start=1):
        if not isinstance(slide, dict):
            raise TypeError(f"slides[{i - 1}] must be a dict")
        for key in ("slide_num", "title", "bullets", "narration"):
            if key not in slide:
                raise ValueError(f"Slide {i} missing '{key}'")
        slide["slide_num"] = i
        if not isinstance(slide["bullets"], list):
            raise TypeError(f"slide {i} 'bullets' must be a list")

    return data


def _soffice_path() -> Path:
    """Resolve LibreOffice's soffice binary; checks env, then common Windows paths, then PATH."""
    env = os.environ.get("LIBREOFFICE_SOFFICE", "").strip()
    if env:
        p = Path(env)
        if p.is_file():
            return p
    for name in ("soffice.exe", "soffice"):
        w = shutil.which(name)
        if w:
            return Path(w)
    for p in (
        Path(r"C:\Program Files\LibreOffice\program\soffice.exe"),
        Path(r"C:\Program Files (x86)\LibreOffice\program\soffice.exe"),
    ):
        if p.is_file():
            return p
    raise FileNotFoundError(
        "LibreOffice (soffice) not found. Install LibreOffice or set LIBREOFFICE_SOFFICE "
        "to the full path of soffice.exe."
    )


def _load_dotenv_stage1() -> None:
    from dotenv import load_dotenv

    load_dotenv(Path(__file__).resolve().parent / ".env")


def _download_image_to_file(url: str, dest: Path) -> None:
    req = urllib.request.Request(
        url,
        headers={"User-Agent": "UdemyCourseGenerator/1.0"},
    )
    with urllib.request.urlopen(req, timeout=90) as resp:
        dest.write_bytes(resp.read())


def resolve_slide_illustrations(
    script: dict[str, Any], output_dir: Path
) -> dict[int, Path]:
    """
    Per slide (1-based index): optional hero image.
    - source_images_dir: relative to output_dir (lesson folder) or absolute; used for bare filenames.
    - slide['image']: https URL, absolute path, path under output_dir, or filename under source_images_dir.
    """
    _load_dotenv_stage1()
    out: dict[int, Path] = {}
    raw_base = script.get("source_images_dir")
    base: Path | None = None
    if isinstance(raw_base, str) and raw_base.strip():
        bp = Path(raw_base.strip())
        base = bp.resolve() if bp.is_absolute() else (output_dir / bp).resolve()
    for slide in script.get("slides") or []:
        try:
            sn = int(slide.get("slide_num", 0))
        except (TypeError, ValueError):
            continue
        if sn < 1:
            continue
        raw_img = slide.get("image")
        if not (isinstance(raw_img, str) and raw_img.strip()):
            continue
        s = raw_img.strip()
        path_dest = output_dir / f"_slide_{sn:02d}_illustration.jpg"
        try:
            if s.lower().startswith("http://") or s.lower().startswith("https://"):
                print(f"  Slide {sn}: downloading image ...")
                _download_image_to_file(s, path_dest)
                if path_dest.is_file() and path_dest.stat().st_size > 0:
                    out[sn] = path_dest
                continue
            p = Path(s)
            if p.is_file():
                out[sn] = p
                continue
            if base is not None:
                bp = (base / s).resolve()
                if bp.is_file():
                    out[sn] = bp
                    continue
            op = (output_dir / s).resolve()
            if op.is_file():
                out[sn] = op
            else:
                print(f"  Warning: slide {sn} image not found: {s}")
        except (urllib.error.URLError, OSError) as e:
            print(f"  Warning: slide {sn} image error: {e}")
    return out


def _build_pptx(
    script: dict[str, Any],
    pptx_path: Path,
    illustrations: dict[int, Path] | None = None,
) -> None:
    """16:9 deck: light background, top accent, title + body, optional right-column photo."""
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt
    from pptx import Presentation

    illustrations = illustrations or {}
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    try:
        blank = prs.slide_layouts[6]
    except IndexError:
        blank = prs.slide_layouts[5]

    for slide in script["slides"]:
        try:
            sn = int(slide["slide_num"])
        except (KeyError, TypeError, ValueError):
            continue
        s = prs.slides.add_slide(blank)
        sw, sh = prs.slide_width, prs.slide_height
        has_pic = sn in illustrations and illustrations[sn].is_file()

        # Order: back to front = background, accent, photo, then text (text on top)
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, sw, sh)
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(248, 250, 252)
        bg.line.fill.background()

        acc = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, sw, Inches(0.12))
        acc.fill.solid()
        acc.fill.fore_color.rgb = RGBColor(30, 64, 175)
        acc.line.fill.background()

        if has_pic and illustrations[sn].is_file():
            try:
                s.shapes.add_picture(
                    str(illustrations[sn].resolve()),
                    Inches(7.65),
                    Inches(0.95),
                    width=Inches(4.9),
                )
            except OSError as e:
                print(f"  Warning: could not add picture to slide {sn}: {e}")
                has_pic = False

        left_w = Inches(6.85) if has_pic else Inches(12.15)
        left0 = Inches(0.55)
        tbox = s.shapes.add_textbox(
            left0, Inches(0.3), left_w, Inches(0.78)
        )
        tf = tbox.text_frame
        tf.clear()
        tf.word_wrap = True
        p0 = tf.paragraphs[0]
        p0.text = str(slide.get("title", ""))
        p0.font.size = Pt(30)
        p0.font.bold = True
        p0.font.color.rgb = RGBColor(15, 23, 42)
        p0.alignment = PP_ALIGN.LEFT

        body = s.shapes.add_textbox(
            left0, Inches(1.1), left_w, Inches(5.6 if has_pic else 6.1)
        )
        btf = body.text_frame
        btf.clear()
        btf.word_wrap = True
        btf.vertical_anchor = MSO_ANCHOR.TOP
        bullets = slide.get("bullets") or []
        for bi, b in enumerate(bullets):
            if bi == 0:
                p = btf.paragraphs[0]
            else:
                p = btf.add_paragraph()
            p.text = str(b)
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(51, 65, 85)
            p.line_spacing = 1.2
            p.space_after = Pt(6)
            p.level = 0

    prs.save(pptx_path)


def generate_slides(script: dict[str, Any], output_dir: Path) -> list[Path]:
    """
    Build a .pptx from script, export to PDF with LibreOffice headless, rasterize each PDF page
    to slide_01.png … slide_N.png with PyMuPDF (subprocess + fitz — one slide per page).
    Direct `soffice --convert-to png` on a deck only yields one PNG on Windows; PDF in the middle fixes that.
    """
    import fitz

    output_dir = output_dir.resolve()
    output_dir.mkdir(parents=True, exist_ok=True)

    pptx_path = output_dir / "deck.pptx"
    pdf_path = output_dir / "deck.pdf"

    print("Building PowerPoint...")
    ill = resolve_slide_illustrations(script, output_dir)
    if ill:
        print(f"  {len(ill)} slide(s) with hero image.")
    _build_pptx(script, pptx_path, ill)
    print(f"  saved {pptx_path.name}")

    soffice = _soffice_path()
    print("Converting to PDF (LibreOffice headless)...")
    r = subprocess.run(
        [
            str(soffice),
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(output_dir),
            str(pptx_path),
        ],
        capture_output=True,
        text=True,
        timeout=180,
        check=False,
    )
    if r.returncode != 0:
        raise RuntimeError(
            f"LibreOffice failed (exit {r.returncode}).\nstdout:\n{r.stdout}\nstderr:\n{r.stderr}"
        )
    if not pdf_path.is_file():
        raise FileNotFoundError(f"Expected {pdf_path} after conversion; got: {r.stdout}")

    print("Rasterizing PDF to PNGs (PyMuPDF)...")
    doc = fitz.open(pdf_path)
    n = len(doc)
    if n != len(script["slides"]):
        doc.close()
        raise RuntimeError(
            f"Slide count mismatch: script has {len(script['slides'])} slides, PDF has {n} pages."
        )

    paths: list[Path] = []
    for i in range(n):
        page = doc[i]
        zoom = 1920.0 / page.rect.width
        pix = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom), alpha=False)
        out = output_dir / f"slide_{i + 1:02d}.png"
        pix.save(out.as_posix())
        paths.append(out)
        print(f"  wrote {out.name}")

    doc.close()
    return paths


def generate_audio(script: dict[str, Any], output_dir: Path) -> list[Path]:
    """
    For each slide, send narration to OpenAI TTS (tts-1-hd) and save MP3s as
    slide_01.mp3 ... slide_N.mp3 in output_dir. Uses the same .env key as the rest of the file.
    """
    output_dir = output_dir.resolve()
    output_dir.mkdir(parents=True, exist_ok=True)
    client = _openai_client()
    n = len(script["slides"])
    paths: list[Path] = []

    for i, slide in enumerate(script["slides"], start=1):
        if "narration" not in slide:
            raise ValueError(f"Slide {i} missing 'narration'")
        text = str(slide["narration"]).strip()
        if not text:
            raise ValueError(f"Slide {i} has empty narration")

        out = output_dir / f"slide_{i:02d}.mp3"
        print(f"TTS ({i}/{n}) {out.name}...")
        r = client.audio.speech.create(
            model="tts-1-hd",
            voice="alloy",
            input=text,
        )
        out.write_bytes(r.read())
        paths.append(out)
        print(f"  saved ({out.stat().st_size} bytes)")

    return paths


# Video filter: pad to 16:9 1920x1080, even sizes for yuv420p, steady fps. Fixes odd heights (e.g. 1920x1081) from PyMuPDF.
_VF_STILL: str = (
    "scale=1920:1080:force_original_aspect_ratio=decrease,"
    "pad=1920:1080:(ow-iw)/2:(oh-ih)/2,format=yuv420p,fps=30"
)


def _ffmpeg_path() -> str:
    w = shutil.which("ffmpeg")
    if not w:
        raise FileNotFoundError(
            "ffmpeg not found on PATH. Install ffmpeg and re-open the terminal, or set PATH."
        )
    return w


def _ffconcat_list_line(p: Path) -> str:
    # Concat demuxer: ' in path must be written as '\''
    s = p.resolve().as_posix()
    s = s.replace("'", r"'\''")
    return f"file '{s}'"


def compose_video(
    slide_pngs: list[Path], audio_mp3s: list[Path], output_path: Path
) -> None:
    """
    For each index: one still image + one MP3 become a short MP4; segments are concat'd (stream copy) into
    output_path. subprocess.run is used to call ffmpeg: same idea as doing it in a terminal, but the script
    can check exit codes and forward stderr if something goes wrong.
    """
    if len(slide_pngs) != len(audio_mp3s):
        raise ValueError(
            f"Image/audio count mismatch: {len(slide_pngs)} PNGs vs {len(audio_mp3s)} MP3s"
        )
    if not slide_pngs:
        raise ValueError("No slides; nothing to compose.")

    output_path = output_path.resolve()
    output_path.parent.mkdir(parents=True, exist_ok=True)
    ffmpeg = _ffmpeg_path()

    tmp = output_path.parent / "_stitch"
    if tmp.is_dir():
        for old in tmp.glob("seg_*.mp4"):
            try:
                old.unlink()
            except OSError:
                pass
    tmp.mkdir(parents=True, exist_ok=True)

    segments: list[Path] = []
    for i, (img, aud) in enumerate(zip(slide_pngs, audio_mp3s, strict=True), start=1):
        seg = tmp / f"seg_{i:02d}.mp4"
        print(f"FFmpeg segment {i}/{len(slide_pngs)}: {img.name} + {aud.name} -> {seg.name}...")
        r = subprocess.run(
            [
                ffmpeg,
                "-y",
                "-hide_banner",
                "-loglevel",
                "error",
                "-loop",
                "1",
                "-i",
                str(img),
                "-i",
                str(aud),
                "-vf",
                _VF_STILL,
                "-c:v",
                "libx264",
                "-tune",
                "stillimage",
                "-c:a",
                "aac",
                "-b:a",
                "192k",
                "-shortest",
                "-movflags",
                "+faststart",
                str(seg),
            ],
            capture_output=True,
            text=True,
            timeout=600,
        )
        if r.returncode != 0 or not seg.is_file():
            err = (r.stderr or r.stdout or "").strip()
            raise RuntimeError(
                f"ffmpeg failed on segment {i} (exit {r.returncode}).\n{err or '(no stderr)'}"
            )
        segments.append(seg)

    list_path = tmp / "file_list.txt"
    list_body = "ffconcat version 1.0\n" + "\n".join(
        _ffconcat_list_line(s) for s in segments
    )
    # UTF-8, no BOM — Required by ffmpeg concat; a Byte Order Mark breaks the first "file" line.
    list_path.write_text(list_body + "\n", encoding="utf-8", newline="\n")

    print("Concatenating segments...")
    r2 = subprocess.run(
        [
            ffmpeg,
            "-y",
            "-hide_banner",
            "-loglevel",
            "error",
            "-f",
            "concat",
            "-safe",
            "0",
            "-i",
            str(list_path),
            "-c",
            "copy",
            str(output_path),
        ],
        capture_output=True,
        text=True,
        timeout=600,
    )
    if r2.returncode != 0 or not output_path.is_file():
        err = (r2.stderr or r2.stdout or "").strip()
        raise RuntimeError(
            f"ffmpeg concat failed (exit {r2.returncode}).\n{err or '(no stderr)'}\n"
            "If stream copy is incompatible, try re-running after deleting _stitch/."
        )

    for s in segments:
        try:
            s.unlink()
        except OSError:
            pass
    try:
        list_path.unlink()
    except OSError:
        pass
    try:
        tmp.rmdir()
    except OSError:
        pass  # e.g. not empty; leave for debugging

    print(f"Wrote {output_path}")


# -----------------------------------------------------------------------------


def main(topic: str) -> None:
    base = Path(__file__).resolve().parent
    out = base / "output"
    out.mkdir(parents=True, exist_ok=True)

    print(f"Topic: {topic!r}\n(1/4) Generating script...")
    script = generate_script(topic)
    (out / "script.json").write_text(
        json.dumps(script, indent=2, ensure_ascii=False) + "\n", encoding="utf-8"
    )
    print(f"  script.json saved, {len(script['slides'])} slide(s)")

    print("\n(2/4) Generating slides (PNG)...")
    pngs = generate_slides(script, out)

    print("\n(3/4) Generating voiceover (MP3)...")
    mp3s = generate_audio(script, out)

    out_mp4 = out / "output.mp4"
    print(f"\n(4/4) Composing {out_mp4.name}...")
    compose_video(pngs, mp3s, out_mp4)
    print(f"\nDone. Play: {out_mp4.resolve()}")


def main_from_json(script_path: Path, output_dir: Path | None = None) -> None:
    """Load an existing script (e.g. from stage2), then slides + TTS + compose only; no new GPT script.
    If output_dir is set, all slide PNGs, MP3s, deck files, and output.mp4 go there (default: stage1/output/)."""
    base = Path(__file__).resolve().parent
    out = (output_dir if output_dir is not None else base / "output").resolve()
    out.mkdir(parents=True, exist_ok=True)
    p = script_path.resolve()
    if not p.is_file():
        raise SystemExit(f"File not found: {p}")
    script: dict[str, Any] = json.loads(p.read_text(encoding="utf-8"))
    if "topic" not in script or "slides" not in script:
        raise SystemExit("script.json must contain 'topic' and 'slides'")
    for i, slide in enumerate(script["slides"], start=1):
        if not isinstance(slide, dict):
            raise SystemExit("Each slide must be an object")
        for k in ("title", "bullets", "narration"):
            if k not in slide:
                raise SystemExit(f"Slide {i} missing '{k}'")
        slide["slide_num"] = i
    (out / "script.json").write_text(
        json.dumps(script, indent=2, ensure_ascii=False) + "\n", encoding="utf-8"
    )
    print(
        f"Using script from: {p}\n(1/3) Generating slides for {len(script['slides'])} slide(s)..."
    )
    pngs = generate_slides(script, out)
    print("\n(2/3) Generating voiceover (MP3)...")
    mp3s = generate_audio(script, out)
    out_mp4 = out / "output.mp4"
    print(f"\n(3/3) Composing {out_mp4.name}...")
    compose_video(pngs, mp3s, out_mp4)
    print(f"\nDone. Play: {out_mp4.resolve()}")


SAMPLE_SCRIPT_FOR_SLIDE_TEST: dict[str, Any] = {
    "topic": "Slide test",
    "slides": [
        {
            "slide_num": 1,
            "title": "First test slide",
            "bullets": ["Bullet A", "Bullet B"],
            "narration": "Not used in generate_slides.",
        },
        {
            "slide_num": 2,
            "title": "Second test slide",
            "bullets": ["Point one", "Point two", "Point three"],
            "narration": "Not used in generate_slides.",
        },
    ],
}

def _print_help() -> None:
    print(
        """\
Usage:
  python make_video.py [TOPIC]
  python make_video.py --test-slides | --test-audio | --test-compose
  python make_video.py --help

  TOPIC        Optional. If omitted, uses a default demo title. Full run writes
               output/script.json, slide_*.png, slide_*.mp3, output/output.mp4.

  --test-slides   Build deck from a built-in 2-slide script (LibreOffice + PyMuPDF).
  --test-audio    Two MP3s from the same built-in script (OpenAI TTS, uses .env).
  --test-compose  Stitch existing slide_*.png + slide_*.mp3 in output/ with ffmpeg.
  --from-json PATH [--out DIR]  Skip GPT script; use script.json for slides+TTS+MP4. Optional --out
               puts all artifacts in DIR (default: stage1/output/). Use for stage4 batch lessons.

  Slides: optional per-slide "image" and top-level "source_images_dir" (see module docstring).
"""
    )


SAMPLE_SCRIPT_FOR_AUDIO_TEST: dict[str, Any] = {
    "topic": "Audio test",
    "slides": [
        {
            "slide_num": 1,
            "title": "A",
            "bullets": ["b"],
            "narration": "This is a short test line for text to speech.",
        },
        {
            "slide_num": 2,
            "title": "B",
            "bullets": ["b"],
            "narration": "Second line, also short, to write slide zero two.",
        },
    ],
}


if __name__ == "__main__":
    if len(sys.argv) >= 2 and sys.argv[1] in ("--help", "-h"):
        _print_help()
        raise SystemExit(0)

    if len(sys.argv) >= 2 and sys.argv[1] == "--test-slides":
        out = Path(__file__).resolve().parent / "output"
        pngs = generate_slides(SAMPLE_SCRIPT_FOR_SLIDE_TEST, out)
        print("OK - slide PNGs:")
        for p in pngs:
            print(" ", p.resolve())
        sys.exit(0)

    if len(sys.argv) >= 2 and sys.argv[1] == "--test-audio":
        out = Path(__file__).resolve().parent / "output"
        mp3s = generate_audio(SAMPLE_SCRIPT_FOR_AUDIO_TEST, out)
        print("OK - audio MP3s:")
        for p in mp3s:
            print(" ", p.resolve())
        sys.exit(0)

    if len(sys.argv) >= 2 and sys.argv[1] == "--from-json":
        if len(sys.argv) < 3:
            raise SystemExit(
                "Usage: python make_video.py --from-json <path\\to\\script.json> [--out <output_dir>]"
            )
        script_path = Path(sys.argv[2])
        out_dir: Path | None = None
        rest = sys.argv[3:]
        i = 0
        while i < len(rest):
            if rest[i] in ("--out", "--output-dir"):
                if i + 1 >= len(rest):
                    raise SystemExit(f"{rest[i]} requires a directory path")
                out_dir = Path(rest[i + 1])
                i += 2
            else:
                raise SystemExit(f"Unknown argument after script path: {rest[i]!r}")
        main_from_json(script_path, out_dir)
        raise SystemExit(0)

    if len(sys.argv) >= 2 and sys.argv[1] == "--test-compose":
        out = Path(__file__).resolve().parent / "output"
        pngs = sorted(out.glob("slide_*.png"))
        mp3s = sorted(out.glob("slide_*.mp3"))
        if not pngs or not mp3s or len(pngs) != len(mp3s):
            raise SystemExit(
                f"Need matching slide_*.png and slide_*.mp3 in {out} "
                f"(found {len(pngs)} PNG, {len(mp3s)} MP3). Run --test-slides and --test-audio first, or a full run."
            )
        out_mp4 = out / "output.mp4"
        compose_video(pngs, mp3s, out_mp4)
        print("OK - composed:", out_mp4.resolve())
        sys.exit(0)

    if len(sys.argv) >= 2 and sys.argv[1].startswith("-"):
        print("Unknown option:", sys.argv[1], file=sys.stderr)
        raise SystemExit(2)

    t = " ".join(sys.argv[1:]) if len(sys.argv) > 1 else "Introduction to Python Variables"
    main(t)
